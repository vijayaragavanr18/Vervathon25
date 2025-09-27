import fitz  # PyMuPDF
from fastapi import UploadFile
from typing import Optional, Dict, List
import os
import logging
import json
import pandas as pd
from sentence_transformers import SentenceTransformer
import numpy as np

class AdvancedDocumentProcessor:
    """Enhanced document processing with PyMuPDF and semantic chunking"""
    
    def __init__(self):
        self.supported_formats = [
            'pdf', 'docx', 'txt', 'md', 'xlsx', 'xls',
            'pptx', 'rtf', 'json', 'csv'
        ]
        self.max_file_size = 50 * 1024 * 1024  # 50MB
        self.temp_dir = "temp_uploads"
        self.chunk_size = 1000  # Characters per chunk
        self.chunk_overlap = 200  # Overlap between chunks
        
        # Initialize sentence transformer for semantic chunking
        self.embedder = None
        self._load_embedder()
        self._ensure_temp_dir()
    
    def _load_embedder(self):
        """Load sentence transformer for semantic processing"""
        try:
            print("📥 Loading sentence transformer for semantic processing...")
            self.embedder = SentenceTransformer('all-MiniLM-L6-v2')
            print("✅ Sentence transformer loaded successfully!")
        except Exception as e:
            print(f"⚠️ Could not load sentence transformer: {e}")
            self.embedder = None
    
    def _ensure_temp_dir(self):
        """Create temp directory if it doesn't exist"""
        if not os.path.exists(self.temp_dir):
            os.makedirs(self.temp_dir)
    
    async def process_upload(self, file: UploadFile) -> Dict:
        """Process uploaded file with advanced extraction"""
        try:
            # Validate file
            validation = await self._validate_file(file)
            if not validation['valid']:
                return validation
            
            # Save temp file
            file_path = await self._save_temp_file(file)
            
            # Process based on file type with PyMuPDF priority
            content_data = await self._extract_content_advanced(file_path, file.filename)
            
            # Clean up temp file
            self._cleanup_file(file_path)
            
            # Create semantic chunks for chat
            chunks = self._create_semantic_chunks(content_data['content'])
            
            # Generate embeddings for retrieval
            embeddings = self._generate_embeddings(chunks)
            
            # Enhanced analysis
            analysis = self._analyze_content_advanced(content_data['content'], file.filename, chunks)
            
            return {
                'success': True,
                'filename': file.filename,
                'content': content_data['content'],
                'structured_content': content_data.get('structured', {}),
                'chunks': chunks,
                'embeddings': embeddings,
                'analysis': analysis,
                'file_type': file.filename.split('.')[-1].lower(),
                'size': len(content_data['content']),
                'pages': content_data.get('pages', 1),
                'extraction_method': content_data.get('method', 'unknown')
            }
            
        except Exception as e:
            logging.error(f"Document processing error: {e}")
            return {
                'success': False,
                'error': str(e),
                'message': 'Failed to process document'
            }
    
    async def _validate_file(self, file: UploadFile) -> Dict:
        """Validate uploaded file"""
        content = await file.read()
        await file.seek(0)
        
        if len(content) > self.max_file_size:
            return {
                'valid': False,
                'error': f'File too large. Maximum size is {self.max_file_size / (1024*1024):.1f}MB'
            }
        
        if not file.filename:
            return {'valid': False, 'error': 'No filename provided'}
        
        extension = file.filename.split('.')[-1].lower()
        if extension not in self.supported_formats:
            return {
                'valid': False,
                'error': f'Unsupported file format. Supported: {", ".join(self.supported_formats)}'
            }
        
        return {'valid': True}
    
    async def _save_temp_file(self, file: UploadFile) -> str:
        """Save uploaded file temporarily"""
        file_path = os.path.join(self.temp_dir, file.filename)
        content = await file.read()
        with open(file_path, 'wb') as f:
            f.write(content)
        return file_path
    
    async def _extract_content_advanced(self, file_path: str, filename: str) -> Dict:
        """Advanced content extraction with PyMuPDF priority"""
        extension = filename.split('.')[-1].lower()
        
        try:
            if extension == 'pdf':
                return await self._extract_pdf_pymupdf(file_path)
            elif extension in ['docx', 'doc']:
                return await self._extract_docx_advanced(file_path)
            elif extension == 'txt':
                return self._extract_txt_advanced(file_path)
            elif extension == 'md':
                return self._extract_markdown_advanced(file_path)
            elif extension == 'json':
                return self._extract_json_advanced(file_path)
            elif extension == 'csv':
                return self._extract_csv_advanced(file_path)
            elif extension in ['xlsx', 'xls']:
                return await self._extract_excel_advanced(file_path)
            elif extension == 'pptx':
                return await self._extract_pptx_advanced(file_path)
            elif extension == 'rtf':
                return await self._extract_rtf_advanced(file_path)
            else:
                return self._extract_txt_advanced(file_path)
        except Exception as e:
            print(f"⚠️ Advanced extraction failed: {e}, falling back to basic extraction")
            return {'content': self._extract_txt_advanced(file_path)['content'], 'method': 'fallback'}
    
    async def _extract_pdf_pymupdf(self, file_path: str) -> Dict:
        """Extract PDF content using PyMuPDF (superior method)"""
        try:
            doc = fitz.open(file_path)
            
            full_text = ""
            structured_data = {
                'pages': [],
                'metadata': {},
                'toc': [],
                'images': []
            }
            
            # Extract metadata
            structured_data['metadata'] = doc.metadata
            
            # Extract table of contents
            toc = doc.get_toc()
            if toc:
                structured_data['toc'] = toc
            
            # Extract content page by page
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                
                # Extract text
                page_text = page.get_text()
                full_text += f"\n--- Page {page_num + 1} ---\n{page_text}\n"
                
                # Extract structured data
                page_data = {
                    'page_number': page_num + 1,
                    'text': page_text,
                    'word_count': len(page_text.split()),
                    'images': len(page.get_images()),
                    'links': len(page.get_links())
                }
                
                structured_data['pages'].append(page_data)
            
            doc.close()
            
            return {
                'content': full_text.strip(),
                'structured': structured_data,
                'pages': len(structured_data['pages']),
                'method': 'pymupdf'
            }
            
        except Exception as e:
            print(f"⚠️ PyMuPDF extraction failed: {e}")
            raise e
    
    async def _extract_docx_advanced(self, file_path: str) -> Dict:
        """Extract DOCX content with structure"""
        try:
            from docx import Document
            doc = Document(file_path)
            
            full_text = ""
            structured_data = {
                'paragraphs': [],
                'tables': [],
                'headers': [],
                'footers': []
            }
            
            # Extract paragraphs with styling
            for para in doc.paragraphs:
                if para.text.strip():
                    full_text += para.text + "\n"
                    structured_data['paragraphs'].append({
                        'text': para.text,
                        'style': para.style.name if para.style else 'Normal'
                    })
            
            # Extract tables
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = []
                    for cell in row.cells:
                        row_data.append(cell.text)
                    table_data.append(row_data)
                structured_data['tables'].append(table_data)
                
                # Add table text to full text
                for row in table_data:
                    full_text += " | ".join(row) + "\n"
            
            return {
                'content': full_text.strip(),
                'structured': structured_data,
                'method': 'python-docx'
            }
            
        except ImportError:
            # Fallback method
            return {'content': 'DOCX file detected but python-docx not available', 'method': 'fallback'}
        except Exception as e:
            print(f"⚠️ DOCX extraction failed: {e}")
            raise e
    
    def _extract_txt_advanced(self, file_path: str) -> Dict:
        """Extract and analyze text file"""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        
        return {
            'content': content,
            'structured': {
                'encoding': 'utf-8',
                'line_count': len(content.split('\n'))
            },
            'method': 'text'
        }
    
    def _extract_markdown_advanced(self, file_path: str) -> Dict:
        """Extract markdown with structure preservation"""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        
        # Extract markdown structure
        import re
        headers = re.findall(r'^(#{1,6})\s+(.+)$', content, re.MULTILINE)
        links = re.findall(r'\[([^\]]+)\]\(([^\)]+)\)', content)
        
        structured_data = {
            'headers': [{'level': len(h[0]), 'text': h[1]} for h in headers],
            'links': [{'text': l[0], 'url': l[1]} for l in links]
        }
        
        # Clean content for chat processing
        clean_content = re.sub(r'#{1,6}\s+', '', content)  # Remove header markers
        clean_content = re.sub(r'\*\*(.*?)\*\*', r'\1', clean_content)  # Remove bold
        clean_content = re.sub(r'\*(.*?)\*', r'\1', clean_content)  # Remove italic
        clean_content = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', clean_content)  # Clean links
        
        return {
            'content': clean_content.strip(),
            'structured': structured_data,
            'method': 'markdown'
        }
    
    def _extract_json_advanced(self, file_path: str) -> Dict:
        """Extract and structure JSON data"""
        with open(file_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        
        # Convert to readable text
        readable_text = self._json_to_readable_text(data)
        
        return {
            'content': readable_text,
            'structured': {'json_data': data},
            'method': 'json'
        }
    
    def _json_to_readable_text(self, data, prefix="") -> str:
        """Convert JSON to conversational text"""
        text_parts = []
        
        if isinstance(data, dict):
            for key, value in data.items():
                if isinstance(value, (dict, list)):
                    text_parts.append(f"{prefix}The {key} contains:")
                    text_parts.append(self._json_to_readable_text(value, prefix + "  "))
                else:
                    text_parts.append(f"{prefix}The {key} is {value}")
        elif isinstance(data, list):
            for i, item in enumerate(data):
                text_parts.append(f"{prefix}Item {i+1}:")
                text_parts.append(self._json_to_readable_text(item, prefix + "  "))
        else:
            text_parts.append(f"{prefix}{data}")
        
        return "\n".join(text_parts)
    
    def _extract_csv_advanced(self, file_path: str) -> Dict:
        """Extract CSV with pandas for better structure"""
        try:
            df = pd.read_csv(file_path, encoding='utf-8', errors='ignore')
            
            # Generate readable summary
            summary_parts = []
            summary_parts.append(f"This dataset contains {len(df)} rows and {len(df.columns)} columns.")
            summary_parts.append(f"Columns: {', '.join(df.columns.tolist())}")
            summary_parts.append("")
            
            # Add sample data
            if len(df) > 0:
                summary_parts.append("Sample data:")
                for idx, row in df.head(5).iterrows():
                    row_desc = []
                    for col, val in row.items():
                        if pd.notna(val):
                            row_desc.append(f"{col}: {val}")
                    summary_parts.append(f"Row {idx + 1}: {', '.join(row_desc)}")
            
            # Add statistics
            if len(df) > 0:
                summary_parts.append("")
                summary_parts.append("Data summary:")
                for col in df.columns:
                    if df[col].dtype in ['int64', 'float64']:
                        summary_parts.append(f"{col}: mean={df[col].mean():.2f}, min={df[col].min()}, max={df[col].max()}")
                    else:
                        unique_count = df[col].nunique()
                        summary_parts.append(f"{col}: {unique_count} unique values")
            
            return {
                'content': "\n".join(summary_parts),
                'structured': {
                    'dataframe_info': df.describe().to_dict(),
                    'columns': df.columns.tolist(),
                    'row_count': len(df),
                    'sample_data': df.head().to_dict('records')
                },
                'method': 'pandas'
            }
            
        except Exception as e:
            return self._extract_csv_basic(file_path)
    
    def _extract_csv_basic(self, file_path: str) -> Dict:
        """Basic CSV extraction fallback"""
        import csv
        text_parts = []
        
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            csv_reader = csv.reader(f)
            headers = next(csv_reader, [])
            
            if headers:
                text_parts.append(f"Columns: {', '.join(headers)}")
                
            for row_num, row in enumerate(csv_reader, 1):
                if row_num > 20:  # Limit for chat processing
                    text_parts.append(f"... and {sum(1 for _ in csv_reader) + 1} more rows")
                    break
                row_text = " | ".join(row)
                text_parts.append(f"Row {row_num}: {row_text}")
        
        return {
            'content': "\n".join(text_parts),
            'method': 'csv_basic'
        }
    
    async def _extract_excel_advanced(self, file_path: str) -> Dict:
        """Advanced Excel extraction with pandas"""
        try:
            # Read all sheets
            excel_file = pd.ExcelFile(file_path)
            all_content = []
            structured_data = {'sheets': {}}
            
            for sheet_name in excel_file.sheet_names:
                df = excel_file.parse(sheet_name)
                
                # Generate readable content for this sheet
                sheet_content = []
                sheet_content.append(f"Sheet: {sheet_name}")
                sheet_content.append(f"Size: {len(df)} rows, {len(df.columns)} columns")
                sheet_content.append(f"Columns: {', '.join(df.columns.tolist())}")
                
                # Sample data
                if len(df) > 0:
                    sheet_content.append("Sample data:")
                    for idx, row in df.head(3).iterrows():
                        row_desc = []
                        for col, val in row.items():
                            if pd.notna(val):
                                row_desc.append(f"{col}: {val}")
                        sheet_content.append(f"  {', '.join(row_desc)}")
                
                all_content.extend(sheet_content)
                all_content.append("")
                
                # Store structured data
                structured_data['sheets'][sheet_name] = {
                    'columns': df.columns.tolist(),
                    'row_count': len(df),
                    'sample_data': df.head().to_dict('records')
                }
            
            return {
                'content': "\n".join(all_content),
                'structured': structured_data,
                'method': 'pandas_excel'
            }
            
        except Exception as e:
            return {'content': 'Excel file detected but pandas not available', 'method': 'fallback'}
    
    async def _extract_pptx_advanced(self, file_path: str) -> Dict:
        """Extract PowerPoint content"""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            
            full_text = ""
            structured_data = {'slides': []}
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_content = []
                slide_content.append(f"Slide {slide_num}:")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(f"  {shape.text}")
                
                slide_text = "\n".join(slide_content)
                full_text += slide_text + "\n\n"
                
                structured_data['slides'].append({
                    'slide_number': slide_num,
                    'content': slide_text
                })
            
            return {
                'content': full_text.strip(),
                'structured': structured_data,
                'method': 'python-pptx'
            }
            
        except ImportError:
            return {'content': 'PowerPoint file detected but python-pptx not available', 'method': 'fallback'}
    
    async def _extract_rtf_advanced(self, file_path: str) -> Dict:
        """Extract RTF content with basic cleanup"""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        
        # Enhanced RTF cleanup
        import re
        content = re.sub(r'\\[a-z]+\d*\s?', '', content)  # Remove RTF commands
        content = re.sub(r'[{}]', '', content)  # Remove braces
        content = re.sub(r'\s+', ' ', content)  # Normalize whitespace
        
        return {
            'content': content.strip(),
            'method': 'rtf'
        }
    
    def _create_semantic_chunks(self, content: str) -> List[Dict]:
        """Create semantic chunks for better chat context"""
        if not content:
            return []
        
        # Split into sentences first
        sentences = self._split_into_sentences(content)
        
        chunks = []
        current_chunk = ""
        current_sentences = []
        
        for sentence in sentences:
            # Check if adding this sentence would exceed chunk size
            if len(current_chunk + sentence) > self.chunk_size and current_chunk:
                # Save current chunk
                chunks.append({
                    'content': current_chunk.strip(),
                    'sentences': current_sentences.copy(),
                    'word_count': len(current_chunk.split()),
                    'char_count': len(current_chunk)
                })
                
                # Start new chunk with overlap
                overlap_content = " ".join(current_sentences[-2:]) if len(current_sentences) >= 2 else ""
                current_chunk = overlap_content + " " + sentence if overlap_content else sentence
                current_sentences = current_sentences[-2:] + [sentence] if len(current_sentences) >= 2 else [sentence]
            else:
                current_chunk += " " + sentence if current_chunk else sentence
                current_sentences.append(sentence)
        
        # Add final chunk
        if current_chunk.strip():
            chunks.append({
                'content': current_chunk.strip(),
                'sentences': current_sentences,
                'word_count': len(current_chunk.split()),
                'char_count': len(current_chunk)
            })
        
        return chunks
    
    def _split_into_sentences(self, text: str) -> List[str]:
        """Split text into sentences using NLTK if available"""
        try:
            import nltk
            from nltk.tokenize import sent_tokenize
            return sent_tokenize(text)
        except:
            # Fallback sentence splitting
            import re
            sentences = re.split(r'[.!?]+', text)
            return [s.strip() for s in sentences if s.strip()]
    
    def _generate_embeddings(self, chunks: List[Dict]) -> List[List[float]]:
        """Generate embeddings for semantic search"""
        if not self.embedder or not chunks:
            return []
        
        try:
            chunk_texts = [chunk['content'] for chunk in chunks]
            embeddings = self.embedder.encode(chunk_texts)
            return embeddings.tolist()
        except Exception as e:
            print(f"⚠️ Embedding generation failed: {e}")
            return []
    
    def _analyze_content_advanced(self, content: str, filename: str, chunks: List[Dict]) -> Dict:
        """Enhanced content analysis"""
        if not content:
            return {}
        
        words = content.split()
        sentences = self._split_into_sentences(content)
        
        # Enhanced analysis
        analysis = {
            'word_count': len(words),
            'sentence_count': len(sentences),
            'character_count': len(content),
            'paragraph_count': len(content.split('\n\n')),
            'chunk_count': len(chunks),
            'estimated_reading_time': max(1, len(words) // 200),
            'language': self._detect_language_advanced(content),
            'content_type': self._detect_content_type_advanced(content, filename),
            'keywords': self._extract_keywords_advanced(content),
            'complexity_score': self._calculate_complexity(content),
            'summary_preview': content[:300] + "..." if len(content) > 300 else content,
            'topics': self._extract_topics(content)
        }
        
        return analysis
    
    def _detect_language_advanced(self, content: str) -> str:
        """Enhanced language detection"""
        # Basic multilingual indicators
        language_indicators = {
            'english': ['the', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for', 'of', 'with', 'this', 'that'],
            'spanish': ['el', 'la', 'y', 'o', 'pero', 'en', 'de', 'con', 'este', 'esta'],
            'french': ['le', 'la', 'et', 'ou', 'mais', 'dans', 'de', 'avec', 'ce', 'cette'],
            'german': ['der', 'die', 'das', 'und', 'oder', 'aber', 'in', 'von', 'mit', 'dieser']
        }
        
        words = content.lower().split()
        if not words:
            return "Unknown"
        
        scores = {}
        for lang, indicators in language_indicators.items():
            score = sum(1 for word in words if word in indicators)
            scores[lang] = score / len(words)
        
        detected_lang = max(scores, key=scores.get)
        confidence = scores[detected_lang]
        
        return detected_lang.capitalize() if confidence > 0.02 else "Unknown"
    
    def _detect_content_type_advanced(self, content: str, filename: str) -> str:
        """Advanced content type detection"""
        content_lower = content.lower()
        
        # Academic patterns
        academic_patterns = [
            'abstract', 'introduction', 'methodology', 'results', 'conclusion', 
            'references', 'bibliography', 'hypothesis', 'literature review',
            'research question', 'data analysis', 'findings'
        ]
        
        # Technical patterns
        technical_patterns = [
            'function', 'class', 'method', 'algorithm', 'database', 'server', 
            'api', 'implementation', 'architecture', 'framework', 'library',
            'code', 'syntax', 'debugging'
        ]
        
        # Business patterns
        business_patterns = [
            'revenue', 'profit', 'market', 'strategy', 'analysis', 'report',
            'quarter', 'financial', 'budget', 'roi', 'kpi', 'metrics'
        ]
        
        # Legal patterns
        legal_patterns = [
            'contract', 'agreement', 'terms', 'conditions', 'clause', 'legal',
            'law', 'regulation', 'compliance', 'liability', 'jurisdiction'
        ]
        
        # Calculate scores
        scores = {
            'Academic Paper': sum(1 for pattern in academic_patterns if pattern in content_lower),
            'Technical Document': sum(1 for pattern in technical_patterns if pattern in content_lower),
            'Business Document': sum(1 for pattern in business_patterns if pattern in content_lower),
            'Legal Document': sum(1 for pattern in legal_patterns if pattern in content_lower)
        }
        
        # File extension hints
        extension = filename.split('.')[-1].lower()
        if extension in ['csv', 'xlsx', 'xls']:
            return "Data/Spreadsheet"
        elif extension == 'json':
            return "Data/Configuration"
        elif extension == 'md':
            return "Documentation"
        elif extension == 'pptx':
            return "Presentation"
        
        # Return highest scoring type or general
        max_score = max(scores.values())
        if max_score > 2:
            return max(scores, key=scores.get)
        
        return "General Document"
    
    def _extract_keywords_advanced(self, content: str) -> List[str]:
        """Advanced keyword extraction with TF-IDF like approach"""
        try:
            import nltk
            from nltk.corpus import stopwords
            from nltk.tokenize import word_tokenize
            
            # Download required NLTK data
            try:
                stop_words = set(stopwords.words('english'))
            except:
                nltk.download('stopwords', quiet=True)
                stop_words = set(stopwords.words('english'))
            
            try:
                tokens = word_tokenize(content.lower())
            except:
                nltk.download('punkt', quiet=True)
                tokens = word_tokenize(content.lower())
            
        except ImportError:
            # Fallback tokenization
            import re
            tokens = re.findall(r'\b\w{3,}\b', content.lower())
            stop_words = set(['the', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for', 'of', 'with', 
                             'a', 'an', 'as', 'are', 'was', 'were', 'been', 'be', 'have', 'has', 'had'])
        
        # Filter tokens
        filtered_tokens = [token for token in tokens if token.isalpha() and len(token) > 2 and token not in stop_words]
        
        # Count frequency
        word_freq = {}
        for token in filtered_tokens:
            word_freq[token] = word_freq.get(token, 0) + 1
        
        # Sort by frequency and return top keywords
        sorted_words = sorted(word_freq.items(), key=lambda x: x[1], reverse=True)
        return [word for word, freq in sorted_words[:15] if freq > 1]
    
    def _calculate_complexity(self, content: str) -> float:
        """Calculate content complexity score"""
        if not content:
            return 0.0
        
        words = content.split()
        sentences = self._split_into_sentences(content)
        
        if not sentences:
            return 0.0
        
        # Average sentence length
        avg_sentence_length = len(words) / len(sentences)
        
        # Average word length
        avg_word_length = sum(len(word) for word in words) / len(words) if words else 0
        
        # Complex word ratio (words with 3+ syllables, approximated by length)
        complex_words = [word for word in words if len(word) > 6]
        complex_word_ratio = len(complex_words) / len(words) if words else 0
        
        # Simple complexity score (0-10 scale)
        complexity = min(10, (avg_sentence_length / 10) + (avg_word_length / 2) + (complex_word_ratio * 5))
        
        return round(complexity, 2)
    
    def _extract_topics(self, content: str) -> List[str]:
        """Extract main topics from content"""
        keywords = self._extract_keywords_advanced(content)
        
        # Group related keywords into topics
        topics = []
        
        # Tech topics
        tech_keywords = [kw for kw in keywords if kw in ['technology', 'software', 'computer', 'data', 'system', 'digital']]
        if tech_keywords:
            topics.append('Technology')
        
        # Business topics
        business_keywords = [kw for kw in keywords if kw in ['business', 'market', 'company', 'financial', 'strategy']]
        if business_keywords:
            topics.append('Business')
        
        # Science topics
        science_keywords = [kw for kw in keywords if kw in ['research', 'study', 'analysis', 'method', 'result']]
        if science_keywords:
            topics.append('Research/Science')
        
        # Education topics
        edu_keywords = [kw for kw in keywords if kw in ['education', 'learning', 'student', 'course', 'knowledge']]
        if edu_keywords:
            topics.append('Education')
        
        return topics[:5]  # Top 5 topics
    
    def search_semantic_chunks(self, query: str, chunks: List[Dict], embeddings: List[List[float]], top_k: int = 3) -> List[Dict]:
        """Search for relevant chunks using semantic similarity"""
        if not self.embedder or not chunks or not embeddings or not query:
            return chunks[:top_k] if chunks else []
        
        try:
            # Generate query embedding
            query_embedding = self.embedder.encode([query])
            
            # Calculate similarities
            similarities = []
            for i, chunk_embedding in enumerate(embeddings):
                # Cosine similarity
                similarity = np.dot(query_embedding[0], chunk_embedding) / (
                    np.linalg.norm(query_embedding[0]) * np.linalg.norm(chunk_embedding)
                )
                similarities.append((i, similarity))
            
            # Sort by similarity and return top chunks
            similarities.sort(key=lambda x: x[1], reverse=True)
            
            relevant_chunks = []
            for i, (chunk_idx, similarity) in enumerate(similarities[:top_k]):
                chunk_copy = chunks[chunk_idx].copy()
                chunk_copy['similarity_score'] = float(similarity)
                chunk_copy['rank'] = i + 1
                relevant_chunks.append(chunk_copy)
            
            return relevant_chunks
            
        except Exception as e:
            print(f"⚠️ Semantic search failed: {e}")
            return chunks[:top_k] if chunks else []
    
    def _cleanup_file(self, file_path: str):
        """Clean up temporary file"""
        try:
            if os.path.exists(file_path):
                os.remove(file_path)
        except Exception as e:
            logging.warning(f"Failed to cleanup temp file {file_path}: {e}")
    
    def get_supported_formats(self) -> List[str]:
        """Get list of supported file formats"""
        return self.supported_formats.copy()