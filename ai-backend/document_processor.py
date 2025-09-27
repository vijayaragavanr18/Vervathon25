from fastapi import UploadFile
from typing import Optional, Dict, List
import os
import logging
import json

class DocumentProcessor:
    """Enhanced document processing for FastAPI backend"""
    
    def __init__(self):
        self.supported_formats = [
            'pdf', 'docx', 'txt', 'md', 'xlsx', 'xls',
            'pptx', 'rtf', 'json', 'csv'
        ]
        self.max_file_size = 50 * 1024 * 1024  # 50MB
        self.temp_dir = "temp_uploads"
        self._ensure_temp_dir()
    
    def _ensure_temp_dir(self):
        """Create temp directory if it doesn't exist"""
        if not os.path.exists(self.temp_dir):
            os.makedirs(self.temp_dir)
    
    async def process_upload(self, file: UploadFile) -> Dict:
        """Process uploaded file and extract content"""
        try:
            # Validate file
            validation = await self._validate_file(file)
            if not validation['valid']:
                return validation
            
            # Save temp file
            file_path = await self._save_temp_file(file)
            
            # Process based on file type
            content = await self._extract_content(file_path, file.filename)
            
            # Clean up temp file
            self._cleanup_file(file_path)
            
            # Analyze content
            analysis = self._analyze_content(content, file.filename)
            
            return {
                'success': True,
                'filename': file.filename,
                'content': content,
                'analysis': analysis,
                'file_type': file.filename.split('.')[-1].lower(),
                'size': len(content)
            }
            
        except Exception as e:
            logging.error(f"Document processing error: {e}")
            return {
                'success': False,
                'error': str(e),
                'message': 'Failed to process document'
            }
    
    async def _validate_file(self, file: UploadFile) -> Dict:
        """Validate uploaded file"""
        # Check file size
        content = await file.read()
        await file.seek(0)  # Reset file pointer
        
        if len(content) > self.max_file_size:
            return {
                'valid': False,
                'error': f'File too large. Maximum size is {self.max_file_size / (1024*1024):.1f}MB'
            }
        
        # Check file extension
        if not file.filename:
            return {
                'valid': False,
                'error': 'No filename provided'
            }
        
        extension = file.filename.split('.')[-1].lower()
        if extension not in self.supported_formats:
            return {
                'valid': False,
                'error': f'Unsupported file format. Supported: {", ".join(self.supported_formats)}'
            }
        
        return {'valid': True}
    
    async def _save_temp_file(self, file: UploadFile) -> str:
        """Save uploaded file temporarily"""
        file_path = os.path.join(self.temp_dir, file.filename)
        
        content = await file.read()
        with open(file_path, 'wb') as f:
            f.write(content)
        
        return file_path
    
    async def _extract_content(self, file_path: str, filename: str) -> str:
        """Extract text content from file"""
        extension = filename.split('.')[-1].lower()
        
        if extension == 'txt':
            return self._extract_txt(file_path)
        elif extension == 'md':
            return self._extract_markdown(file_path)
        elif extension == 'json':
            return self._extract_json(file_path)
        elif extension == 'csv':
            return self._extract_csv(file_path)
        elif extension == 'pdf':
            return await self._extract_pdf(file_path)
        elif extension in ['docx', 'doc']:
            return await self._extract_docx(file_path)
        elif extension in ['xlsx', 'xls']:
            return await self._extract_excel(file_path)
        elif extension == 'pptx':
            return await self._extract_pptx(file_path)
        elif extension == 'rtf':
            return await self._extract_rtf(file_path)
        else:
            # Fallback to text extraction
            return self._extract_txt(file_path)
    
    def _extract_txt(self, file_path: str) -> str:
        """Extract content from text file"""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
    
    def _extract_markdown(self, file_path: str) -> str:
        """Extract content from markdown file"""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
            # Remove markdown formatting for cleaner text
            import re
            content = re.sub(r'#{1,6}\s+', '', content)  # Headers
            content = re.sub(r'\*\*(.*?)\*\*', r'\1', content)  # Bold
            content = re.sub(r'\*(.*?)\*', r'\1', content)  # Italic
            content = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', content)  # Links
            return content
    
    def _extract_json(self, file_path: str) -> str:
        """Extract content from JSON file"""
        with open(file_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
            # Convert JSON to readable text
            return self._json_to_text(data)
    
    def _json_to_text(self, data, prefix="") -> str:
        """Convert JSON data to readable text"""
        text_parts = []
        
        if isinstance(data, dict):
            for key, value in data.items():
                if isinstance(value, (dict, list)):
                    text_parts.append(f"{prefix}{key}:")
                    text_parts.append(self._json_to_text(value, prefix + "  "))
                else:
                    text_parts.append(f"{prefix}{key}: {value}")
        elif isinstance(data, list):
            for i, item in enumerate(data):
                text_parts.append(f"{prefix}Item {i+1}:")
                text_parts.append(self._json_to_text(item, prefix + "  "))
        else:
            text_parts.append(f"{prefix}{data}")
        
        return "\n".join(text_parts)
    
    def _extract_csv(self, file_path: str) -> str:
        """Extract content from CSV file"""
        import csv
        text_parts = []
        
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            csv_reader = csv.reader(f)
            headers = next(csv_reader, [])
            
            if headers:
                text_parts.append(f"Columns: {', '.join(headers)}")
                text_parts.append("")
            
            for row_num, row in enumerate(csv_reader, 1):
                if row_num > 100:  # Limit to first 100 rows
                    text_parts.append(f"... and {sum(1 for _ in csv_reader)} more rows")
                    break
                
                row_text = " | ".join(row)
                text_parts.append(f"Row {row_num}: {row_text}")
        
        return "\n".join(text_parts)
    
    async def _extract_pdf(self, file_path: str) -> str:
        """Extract content from PDF file (fallback method)"""
        try:
            # Try using PyPDF2 if available
            import PyPDF2
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                text_parts = []
                for page in reader.pages:
                    text_parts.append(page.extract_text())
                return "\n".join(text_parts)
        except ImportError:
            # Fallback: read as binary and try to extract readable text
            with open(file_path, 'rb') as f:
                content = f.read()
                # Simple text extraction (not ideal but works for some PDFs)
                text = content.decode('latin-1', errors='ignore')
                # Clean up the text
                import re
                text = re.sub(r'[^\x20-\x7E\n]', '', text)
                return text
    
    async def _extract_docx(self, file_path: str) -> str:
        """Extract content from DOCX file (fallback method)"""
        try:
            # Try using python-docx if available
            from docx import Document
            doc = Document(file_path)
            text_parts = []
            for paragraph in doc.paragraphs:
                text_parts.append(paragraph.text)
            return "\n".join(text_parts)
        except ImportError:
            # Fallback: extract from docx as zip
            import zipfile
            import xml.etree.ElementTree as ET
            
            with zipfile.ZipFile(file_path, 'r') as docx:
                content = docx.read('word/document.xml')
                root = ET.fromstring(content)
                
                # Extract text from XML
                text_parts = []
                for elem in root.iter():
                    if elem.text:
                        text_parts.append(elem.text)
                
                return " ".join(text_parts)
    
    async def _extract_excel(self, file_path: str) -> str:
        """Extract content from Excel file (fallback method)"""
        try:
            # Try using openpyxl if available
            from openpyxl import load_workbook
            wb = load_workbook(file_path, data_only=True)
            text_parts = []
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text_parts.append(f"Sheet: {sheet_name}")
                
                for row in sheet.iter_rows(max_row=50, values_only=True):
                    row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                    if row_text.strip():
                        text_parts.append(row_text)
                
                text_parts.append("")
            
            return "\n".join(text_parts)
        except ImportError:
            return "Excel file detected but openpyxl not available for processing."
    
    async def _extract_pptx(self, file_path: str) -> str:
        """Extract content from PowerPoint file"""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text_parts = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                text_parts.append(f"Slide {slide_num}:")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_parts.append(shape.text)
                
                text_parts.append("")
            
            return "\n".join(text_parts)
        except ImportError:
            return "PowerPoint file detected but python-pptx not available for processing."
    
    async def _extract_rtf(self, file_path: str) -> str:
        """Extract content from RTF file (basic)"""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
            # Basic RTF cleanup
            import re
            content = re.sub(r'\\[a-z]+\d*', '', content)  # Remove RTF commands
            content = re.sub(r'[{}]', '', content)  # Remove braces
            content = re.sub(r'\s+', ' ', content)  # Normalize whitespace
            return content.strip()
    
    def _analyze_content(self, content: str, filename: str) -> Dict:
        """Analyze extracted content"""
        words = content.split()
        sentences = content.split('.')
        
        # Detect language (basic)
        language = self._detect_language(content)
        
        # Detect content type
        content_type = self._detect_content_type(content, filename)
        
        # Extract key topics (basic keyword extraction)
        keywords = self._extract_keywords(content)
        
        return {
            'word_count': len(words),
            'sentence_count': len(sentences),
            'character_count': len(content),
            'estimated_reading_time': max(1, len(words) // 200),  # ~200 WPM
            'language': language,
            'content_type': content_type,
            'keywords': keywords[:10],  # Top 10 keywords
            'summary_preview': content[:200] + "..." if len(content) > 200 else content
        }
    
    def _detect_language(self, content: str) -> str:
        """Simple language detection"""
        # Basic language indicators
        english_indicators = ['the', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for', 'of', 'with']
        
        words = content.lower().split()
        english_count = sum(1 for word in words if word in english_indicators)
        
        if len(words) > 0 and english_count / len(words) > 0.05:
            return "English"
        else:
            return "Unknown"
    
    def _detect_content_type(self, content: str, filename: str) -> str:
        """Detect the type of content"""
        content_lower = content.lower()
        
        # Academic indicators
        academic_keywords = ['abstract', 'introduction', 'methodology', 'conclusion', 'references', 'bibliography']
        if any(keyword in content_lower for keyword in academic_keywords):
            return "Academic Paper"
        
        # Technical indicators
        tech_keywords = ['function', 'class', 'method', 'algorithm', 'database', 'server', 'api']
        if any(keyword in content_lower for keyword in tech_keywords):
            return "Technical Document"
        
        # Business indicators
        business_keywords = ['revenue', 'profit', 'market', 'strategy', 'analysis', 'report', 'quarter']
        if any(keyword in content_lower for keyword in business_keywords):
            return "Business Document"
        
        # Educational indicators
        edu_keywords = ['chapter', 'lesson', 'exercise', 'question', 'answer', 'quiz', 'test']
        if any(keyword in content_lower for keyword in edu_keywords):
            return "Educational Content"
        
        # Based on file extension
        extension = filename.split('.')[-1].lower()
        if extension in ['csv', 'xlsx', 'xls']:
            return "Data/Spreadsheet"
        elif extension == 'json':
            return "Data/Configuration"
        elif extension == 'md':
            return "Documentation"
        
        return "General Document"
    
    def _extract_keywords(self, content: str) -> List[str]:
        """Extract key terms from content"""
        # Simple keyword extraction
        words = content.lower().split()
        
        # Remove common stop words
        stop_words = set(['the', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for', 'of', 'with', 
                         'a', 'an', 'as', 'are', 'was', 'were', 'been', 'be', 'have', 'has', 'had',
                         'do', 'does', 'did', 'will', 'would', 'could', 'should', 'may', 'might',
                         'must', 'can', 'this', 'that', 'these', 'those', 'i', 'you', 'he', 'she',
                         'it', 'we', 'they', 'me', 'him', 'her', 'us', 'them', 'my', 'your', 'his',
                         'her', 'its', 'our', 'their'])
        
        # Filter and count words
        word_freq = {}
        for word in words:
            # Clean word
            import re
            word = re.sub(r'[^\w]', '', word)
            
            if len(word) > 3 and word not in stop_words:
                word_freq[word] = word_freq.get(word, 0) + 1
        
        # Sort by frequency
        sorted_words = sorted(word_freq.items(), key=lambda x: x[1], reverse=True)
        
        return [word for word, freq in sorted_words[:20]]
    
    def _cleanup_file(self, file_path: str):
        """Clean up temporary file"""
        try:
            if os.path.exists(file_path):
                os.remove(file_path)
        except Exception as e:
            logging.warning(f"Failed to cleanup temp file {file_path}: {e}")
    
    def get_supported_formats(self) -> List[str]:
        """Get list of supported file formats"""
        return self.supported_formats.copy()